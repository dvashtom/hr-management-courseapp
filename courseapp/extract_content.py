"""
סקריפט לחילוץ תוכן מקבצי PPTX ו-DOCX והמרתם לקבצי Markdown.
סורק את כל תיקיות הפרויקט באופן רקורסיבי.
"""

import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    from docx import Document
except ImportError:
    print("Error: python-docx not installed. Run: pip install python-docx")
    sys.exit(1)


def extract_pptx_to_markdown(pptx_path):
    """חילוץ טקסט מקובץ PPTX והמרתו ל-Markdown"""
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"  ⚠️ שגיאה בפתיחת {pptx_path}: {e}")
        return None

    md_lines = []
    filename = Path(pptx_path).stem
    md_lines.append(f"# {filename}\n")

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_title = None
        slide_content = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue

                    # First text found that looks like a title
                    if slide_title is None and len(text) < 100:
                        slide_title = text
                    else:
                        slide_content.append(text)

        if slide_title or slide_content:
            md_lines.append(f"\n## שקף {slide_num}" + (f": {slide_title}" if slide_title else ""))
            md_lines.append("")
            for line in slide_content:
                # If it looks like a bullet point
                if line.startswith(('-', '•', '*', '–')):
                    md_lines.append(f"- {line.lstrip('-•*– ')}")
                else:
                    md_lines.append(f"- {line}")
            md_lines.append("")

    return "\n".join(md_lines) if len(md_lines) > 2 else None


def extract_docx_to_markdown(docx_path):
    """חילוץ טקסט מקובץ DOCX והמרתו ל-Markdown"""
    try:
        doc = Document(docx_path)
    except Exception as e:
        print(f"  ⚠️ שגיאה בפתיחת {docx_path}: {e}")
        return None

    md_lines = []
    filename = Path(docx_path).stem
    md_lines.append(f"# {filename}\n")

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            md_lines.append("")
            continue

        style_name = paragraph.style.name.lower() if paragraph.style else ""

        # Detect headings
        if 'heading 1' in style_name or 'title' in style_name:
            md_lines.append(f"\n## {text}\n")
        elif 'heading 2' in style_name:
            md_lines.append(f"\n### {text}\n")
        elif 'heading 3' in style_name:
            md_lines.append(f"\n#### {text}\n")
        elif 'list' in style_name or text.startswith(('-', '•', '*', '–')):
            md_lines.append(f"- {text.lstrip('-•*– ')}")
        elif paragraph.runs and paragraph.runs[0].bold:
            md_lines.append(f"\n**{text}**\n")
        else:
            md_lines.append(text)

    return "\n".join(md_lines) if len(md_lines) > 2 else None


def process_directory(base_dir):
    """סריקת התיקייה וחילוץ תוכן מכל הקבצים"""
    base_path = Path(base_dir)
    processed = 0
    skipped = 0
    errors = 0

    print(f"\n🔍 סורק את התיקייה: {base_path}")
    print("=" * 60)

    for root, dirs, files in os.walk(base_path):
        # Skip courseapp directory itself
        if 'courseapp' in root:
            continue

        for file in sorted(files):
            file_path = Path(root) / file
            ext = file_path.suffix.lower()

            if ext not in ('.pptx', '.ppsx', '.docx'):
                continue

            # Skip temp files
            if file.startswith('~'):
                continue

            md_path = file_path.with_suffix('.md')

            # Check if already processed
            if md_path.exists():
                print(f"  ⏭️  כבר קיים: {md_path.name}")
                skipped += 1
                continue

            print(f"  📄 מעבד: {file_path.relative_to(base_path)}")

            if ext in ('.pptx', '.ppsx'):
                content = extract_pptx_to_markdown(str(file_path))
            elif ext == '.docx':
                content = extract_docx_to_markdown(str(file_path))
            else:
                continue

            if content:
                md_path.write_text(content, encoding='utf-8')
                print(f"     ✅ נוצר: {md_path.name}")
                processed += 1
            else:
                print(f"     ⚠️ לא נמצא תוכן")
                errors += 1

    print("\n" + "=" * 60)
    print(f"📊 סיכום:")
    print(f"   ✅ עובדו בהצלחה: {processed}")
    print(f"   ⏭️  כבר קיימים: {skipped}")
    print(f"   ⚠️  שגיאות: {errors}")
    print("=" * 60)


if __name__ == "__main__":
    # The base directory is the parent of courseapp
    script_dir = Path(__file__).parent
    base_dir = script_dir.parent

    print("🎓 סקריפט חילוץ תוכן - ניהול משאבי אנוש")
    print(f"📁 תיקיית בסיס: {base_dir}")

    process_directory(base_dir)
    print("\n✨ הסתיים! כעת ניתן להריץ את האפליקציה:")
    print("   streamlit run app.py")